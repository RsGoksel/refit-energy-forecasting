"""
Generate the EBT 629E project presentation as PPTX.
Uses python-pptx to build slides with embedded figures from results/.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
RESULTS = os.path.join(BASE, "results")
OUTPUT = os.path.join(BASE, "EBT629E_Project_Presentation.pptx")

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x00, 0x33, 0x66)
GREY = RGBColor(0x4D, 0x4D, 0x4D)


def add_title_slide(title, subtitle, members):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(1.0))
    p = tx2.text_frame.paragraphs[0]
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = GREY
    # Members
    tx3 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(1.5))
    tf = tx3.text_frame
    for i, m in enumerate(members):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = m
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = GREY
    return slide


def add_title_content_slide(title, bullets):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Header
    th = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = th.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # Bullets
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = tb.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(20)
        p.font.color.rgb = GREY
        p.space_after = Pt(8)
    return slide


def add_title_image_slide(title, image_filename, caption=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    th = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = th.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # Image, centered, height up to 5.4 in
    img_path = os.path.join(RESULTS, image_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.7), Inches(1.3), height=Inches(5.4))
    if caption:
        tc = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
        p = tc.text_frame.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = GREY
    return slide


def add_table_slide(title, headers, rows):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    th = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = th.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    cols = len(headers)
    n_rows = len(rows) + 1
    left, top = Inches(1.5), Inches(2.0)
    width, height = Inches(10.0), Inches(3.0)
    tbl = slide.shapes.add_table(n_rows, cols, left, top, width, height).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(16)
    for r_idx, row in enumerate(rows):
        for c_idx, v in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
    return slide


# ============ BUILD SLIDES ============

add_title_slide(
    "Household Energy Consumption Forecasting",
    "A Comparative Study of ARIMA, LSTM, and FB-Prophet on the REFIT Dataset",
    [
        "Nurullah Yildirim (301252004)",
        "Kadir Goksel Gunduz (301241077)",
        "Furkan Cinar (301212001)",
        "",
        "EBT 629E, Artificial Intelligence, ITU, May 2026",
    ],
)

add_title_content_slide(
    "Why Forecast Household Energy?",
    [
        "Smart meter roll-out makes per-house forecasting feasible.",
        "Useful for demand response, PV self-consumption, battery scheduling.",
        "Single-house signals are noisy, idiosyncratic, hard to forecast.",
        "Question: ARIMA vs LSTM vs Prophet, which works best at this scale?",
    ],
)

add_title_image_slide(
    "REFIT Dataset, House 2 Appliance Correlations",
    "correlation_heatmap.png",
    "20 UK households, 8-second sampling, October 2013 to June 2015.",
)

add_title_image_slide(
    "Exploratory Data Analysis",
    "data_analysis.png",
    "Evening peaks 16:00 to 20:00, weekend uplift, right-skewed daily distribution.",
)

add_title_image_slide(
    "End-to-End Pipeline",
    "pipeline_diagram.png",
    "From 5.7M raw rows to 122 test days through daily resample and 7-day smoothing.",
)

add_title_content_slide(
    "ARIMA Approach",
    [
        "ADF stationarity test: p-value 0.94, so d = 1.",
        "Grid search over (p, q) in [0,2], select by lowest AIC.",
        "Best model: ARIMA(2, 1, 2), AIC = 4547.6.",
        "Deployed under rolling one-step-ahead forecast.",
    ],
)

add_title_image_slide(
    "LSTM Architecture",
    "lstm_architecture.png",
    "3 stacked LSTM layers + dropout, 14-day input window, 128,929 trainable params.",
)

add_title_image_slide(
    "FB-Prophet Decomposition",
    "prophet_components.png",
    "Weekly + yearly seasonality, multiplicative mode, default changepoint scale.",
)

add_title_image_slide(
    "KEY SLIDE: Rolling vs Multi-step Protocol",
    "rolling_forecast.png",
    "ARIMA / LSTM see yesterday's actual. Prophet must commit to 122 days at once.",
)

add_table_slide(
    "Final Performance",
    ["Model", "MSE", "RMSE (W)", "MAE (W)", "R2"],
    [
        ["ARIMA(2,1,2)", "2,370.01", "48.68", "32.75", "0.876"],
        ["LSTM (3-layer)", "15,436.04", "124.24", "93.59", "0.193"],
        ["FB-Prophet", "18,851.52", "137.30", "111.70", "0.015"],
    ],
)

add_title_image_slide(
    "Predicted vs Actual",
    "model_comparison.png",
    "ARIMA tracks closely, LSTM smooths peaks, Prophet captures only the baseline.",
)

add_title_image_slide(
    "Metrics Side by Side",
    "metrics_comparison.png",
)

add_title_image_slide(
    "Development Iterations",
    "iteration_timeline.png",
    "Daily resample + 7-day smoothing + rolling forecast lifted R2 from -0.13 to 0.876.",
)

add_title_content_slide(
    "Discussion",
    [
        "ARIMA wins: strong lag-1 autocorrelation + rolling refit.",
        "LSTM underperforms: only 476 training sequences, data-scale bound.",
        "Prophet looks worst: unfair multi-step protocol, no test access.",
        "Hybrid models (CNN-LSTM-Transformer, STL-Prophet-LSTM) are the natural next step.",
    ],
)

add_title_content_slide(
    "Conclusion and Future Work",
    [
        "Rolling protocol is the single most impactful design choice.",
        "Classical statistics stay competitive at single-household daily granularity.",
        "Future: weather + calendar features via SARIMAX / multivariate LSTM.",
        "Future: cross-household transfer learning across all 20 REFIT homes.",
        "Future: rolling refit Prophet for a fair multi-step comparison.",
    ],
)

prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(OUTPUT)/1024:.0f} KB")
